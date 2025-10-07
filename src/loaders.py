# src/loaders.py
from __future__ import annotations
from pathlib import Path
from typing import List, Union
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

import os, csv
import chardet
from pypdf import PdfReader
from pptx import Presentation
import openpyxl
from src.config import settings

# [ADD OCR] 의존성
# >>> FIX: pdf2image / pytesseract가 없을 경우를 대비해 안전 가드 추가
try:
    from pdf2image import convert_from_path            # poppler-utils 필요
    _PDF2IMAGE_OK = True   # ← 의존성 확인 플래그
except Exception:
    convert_from_path = None
    _PDF2IMAGE_OK = False  # ← 없으면 False

try:
    import pytesseract                                 # tesseract-ocr / kor / eng 필요
    _TESSERACT_OK = True
except Exception:
    pytesseract = None
    _TESSERACT_OK = False

try:
    import cv2
    import numpy as np
    _CV_OK = True
except Exception:
    cv2 = None
    np = None
    _CV_OK = False
from PIL import Image
import re

# [ADD OCR] 기본 OCR 옵션
DEFAULT_OCR_LANG = "kor+eng"
DEFAULT_OCR_DPI = 300
AUTO_OCR_MIN_CHARS = 80   # auto 경로에서 텍스트가 이 값보다 작으면 OCR 폴백

# === Universal heading detection (범용 헤더 스플리터) ===
HEADING_RES = [
    # 법/규정: 제7조, 제│7 조 … (세로바/기호 허용)
    (re.compile(
        r'^\s*(?:[•\-\u25BA\u25CF\u2022▶]?\s*)?(?P<title>제\s*[\|\│\.\-\u2502\u2503]?\s*(?P<num>\d{1,3})\s*조[^\n]*?)\s*$',
        re.MULTILINE
     ), 1, "article"),
    # 1. 제목
    (re.compile(r'^\s*(?P<title>\d{1,3}\.\s+[^\n]{2,})\s*$', re.MULTILINE), 1, "h1"),
    # 1.1. …
    (re.compile(r'^\s*(?P<title>\d{1,3}(?:\.\d{1,3}){1,}\.?\s+[^\n]{2,})\s*$', re.MULTILINE), 2, "h2"),
    # 로마숫자
    (re.compile(r'^\s*(?P<title>[IVXLCDM]+\.\s+[^\n]{2,})\s*$', re.IGNORECASE | re.MULTILINE), 1, "roman"),
    # 가. 나. 다.
    (re.compile(r'^\s*(?P<title>[가-힣]\.\s+[^\n]{2,})\s*$', re.MULTILINE), 2, "alpha_ko"),
    # (1) (2) …
    (re.compile(r'^\s*(?P<title>\(\d+\)\s+[^\n]{2,})\s*$', re.MULTILINE), 3, "paren"),
]

def _split_sections(text: str) -> list[dict]:
    if not text or len(text) < 30:
        return []
    idxs = []
    for pat, level, kind in HEADING_RES:
        for m in pat.finditer(text):
            title = (m.group("title") or m.group(0)).strip()
            idxs.append((m.start(), m.end(), title, level, kind))

    if not idxs:
        # Fallback: 줄 시작이 아니어도 '제6조'가 보이면 그 지점부터 섹션으로 간주
        ANY_ART_RE = re.compile(r'제\s*(\d{1,3})\s*조[^\n]*')
        idxs2 = [(m.start(), m.end(), m.group(0).strip(), 1, "article") for m in ANY_ART_RE.finditer(text)]
        if not idxs2:
            return []
        idxs = sorted(idxs2, key=lambda x: x[0])

    # 시작 위치 기준으로 제일 '강한' 헤더 하나만 남기기 (level 낮을수록 상위)
    by_start = {}
    for s, e, title, level, kind in idxs:
        if s not in by_start or level < by_start[s][3]:
            by_start[s] = (s, e, title, level, kind)
    idxs = sorted(by_start.values(), key=lambda x: x[0])

    sections = []
    for i, (s, e, title, level, kind) in enumerate(idxs):
        body_start = e
        body_end = idxs[i+1][0] if i + 1 < len(idxs) else len(text)
        body = text[body_start:body_end].strip()
        if len(body) < 40:
            continue
        sections.append({"title": title, "level": level, "kind": kind, "body": body})
    return sections

def _build_section_docs_from_text(text: str, path: Path) -> List[Document]:
    """split_sections 결과를 LangChain Document로 변환 + article_no 메타 자동 부여"""
    # 조문 표기 정규화 먼저
    text = _normalize_articles(text)

    secs = _split_sections(text)[:400]
    docs: List[Document] = []
    for i, sec in enumerate(secs):
        md = {
            "source": str(path),
            "kind": "section",
            "section_title": sec["title"],
            "section_level": sec["level"],
            "section_index": i
        }
        m = re.search(r'제\s*(\d{1,3})\s*조', sec["title"])
        if m:
            try:
                md["article_no"] = int(m.group(1))
            except Exception:
                pass
        # 제목은 prefix에서 넣어주므로 본문만 저장
        docs.append(Document(page_content=sec["body"], metadata=md))
    return docs

_OCR_SEP_RE = re.compile(r"[│|¦┃┆┇┊┋丨ㅣ]")  # 세로바/유사문자

def _norm_brackets(s: str) -> str:
    repl = {"［":"[", "］":"]", "【":"[", "】":"]", "〔":"(", "〕":")", "「":"[", "」":"]", "『":"[", "』":"]"}
    for a, b in repl.items():
        s = s.replace(a, b)
    return s  # ← 반드시 반환

def _normalize_articles(s: str) -> str:
    s = _norm_brackets(s)
    s = _OCR_SEP_RE.sub(" ", s)  # │, | 등 -> 공백
    # 제 (세로바/점/하이픈 허용) 숫자 조  ->  제{숫자}조
    pat = re.compile(r"제\s*[\|\│\.\-\u2502\u2503]?\s*([0-9]{1,3})\s*조")
    return pat.sub(lambda m: f"제{int(m.group(1))}조", s)


# -----------------------------
# 공통: 청킹
# -----------------------------
def _chunk(docs: List[Document], chunk_size=800, chunk_overlap=120) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    out: List[Document] = []
    for d in docs:
        text = d.page_content or ""
        prefix = ""
        md = d.metadata or {}

        if md.get("kind") == "section":
            st = md.get("section_title") or ""
            src_bn = Path(str(md.get("source",""))).name
            if st:
                prefix = f"[SOURCE:{src_bn}] {st}\n"
            else:
                prefix = f"[SOURCE:{src_bn}]\n"

        parts = splitter.split_text(text)
        for i, p in enumerate(parts):
            new_md = dict(md)
            new_md["chunk"] = i
            out.append(Document(page_content=(prefix + p) if prefix else p, metadata=new_md))
    return out

# -----------------------------
# TXT
# -----------------------------
def load_txt(path: Path) -> List[Document]:
    raw = path.read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    text = raw.decode(enc, errors="ignore")
    return [Document(page_content=text, metadata={"source": str(path), "type": "txt"})]

# -----------------------------
# PDF (기본 텍스트 파서)
# -----------------------------
def load_pdf(path: Path) -> List[Document]:
    reader = PdfReader(str(path))
    docs: List[Document] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            docs.append(Document(page_content=text, metadata={
                "source": str(path), "type": "pdf", "page": i
            }))
    return docs

# -----------------------------
# PPTX
# -----------------------------
def load_pptx(path: Path) -> List[Document]:
    prs = Presentation(str(path))
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and shp.text_frame:
                for p in shp.text_frame.paragraphs:
                    if p.runs:
                        lines.append("".join(r.text or "" for r in p.runs))
                    else:
                        lines.append(p.text or "")
            if getattr(shp, "has_table", False):
                tbl = shp.table
                for row in tbl.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    lines.append("\t".join(cells))
        text = "\n".join([ln for ln in lines if ln.strip()])
        if text.strip():
            docs.append(Document(page_content=text, metadata={
                "source": str(path), "type": "pptx", "slide": i
            }))
    return docs

# -----------------------------
# XLSX
# -----------------------------
def load_xlsx(path: Path) -> List[Document]:
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    docs: List[Document] = []
    for ws in wb.worksheets:
        lines = []
        for row in ws.iter_rows(values_only=True):
            line = "\t".join("" if v is None else str(v) for v in row)
            if line.strip():
                lines.append(line)
        text = "\n".join(lines)
        if text.strip():
            docs.append(Document(page_content=text, metadata={
                "source": str(path), "type": "xlsx", "sheet": ws.title
            }))
    return docs

# -----------------------------
# CSV (제목/요약 메타 청크 추가)
# -----------------------------
def load_csv(path: Path) -> List[Document]:
    raw = path.read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    text = raw.decode(enc, errors="ignore")

    docs: List[Document] = []
    rows_norm: list[dict] = []
    lines: list[str] = []

    # DictReader 우선 (헤더가 있을 때)
    rdr = csv.DictReader(text.splitlines())
    if rdr.fieldnames:
        for r in rdr:
            r = { (k or "").strip(): (v or "").strip() for k, v in r.items() }
            rows_norm.append(r)
            k = r.get("항목") or r.get("항목명") or r.get("key") or ""
            v = r.get("내용") or r.get("세부내용") or r.get("value") or ""
            line = f"{k}: {v}" if k else v
            if line.strip():
                lines.append(line)
    else:
        # 헤더가 없으면 일반 reader로 라인 합치기
        rdr2 = csv.reader(text.splitlines())
        for r in rdr2:
            line = ",".join("" if v is None else str(v) for v in r).strip()
            if line:
                lines.append(line)

    if lines:
        docs.append(Document(
            page_content="\n".join(lines),
            metadata={"source": str(path), "type": "csv"}
        ))

    # === TITLE ===
    title = None
    for key in ("제목", "주제", "문서명", "title", "subject"):
        for r in rows_norm:
            if r.get(key):
                title = r[key]
                break
        if title:
            break
    if not title:
        # 항목/항목명에 적절한 값이 있으면 사용
        for r in rows_norm:
            maybe = r.get("항목") or r.get("항목명")
            if maybe and len(maybe) >= 2:
                title = maybe
                break
    if not title:
        title = path.stem

    docs.append(Document(
        page_content=f"[TITLE] {title}",
        metadata={"source": str(path), "kind": "title"}
    ))

    # === SUMMARY === (상위 몇 개를 key=value로)
    pairs = []
    for r in rows_norm[:6]:
        k = (r.get("항목") or r.get("항목명") or "").strip()
        v = (r.get("내용") or r.get("세부내용") or "").strip()
        if k and v:
            pairs.append(f"{k}={v}")
    if pairs:
        docs.append(Document(
            page_content=f"[SUMMARY] {path.name} 개요: " + "; ".join(pairs),
            metadata={"source": str(path), "kind": "summary"}
        ))
    return docs

# -----------------------------
# PDF (표 친화 자동 로더)
# -----------------------------
def _norm(s: str | None) -> str:
    return " ".join((s or "").split())

def _table_to_sentences(table: list[list[str]]) -> list[str]:
    rows = [[_norm(c) for c in r] for r in table if any(_norm(c) for c in r)]
    if not rows:
        return []
    headers = rows[0]
    out = []
    for r in rows[1:]:
        pairs = []
        for h, v in zip(headers, r):
            if _norm(h) and _norm(v):
                pairs.append(f"{h}={v}")
        if pairs:
            out.append(", ".join(pairs))
    return out

def _dedup(docs: List[Document]) -> List[Document]:
    seen = set()
    out: List[Document] = []
    for d in docs:
        key = d.page_content.strip()
        if key and key not in seen:
            seen.add(key)
            out.append(d)
    return out

# [ADD] 사람이 읽을만한지 검사
def _is_readable(s: str) -> bool:
    s = (s or "").strip()
    if not s: return False
    good = len(re.findall(r"[ㄱ-ㅎ가-힣A-Za-z0-9\s\-\_\.\(\)\[\]]", s))
    return good / max(1, len(s)) >= 0.6

def _extract_pdf_title(path: Path) -> str | None:
    try:
        meta = PdfReader(str(path)).metadata or {}
        t = getattr(meta, "title", None)
        if isinstance(meta, dict):
            t = t or meta.get("/Title") or meta.get("Title")
        if t:
            t = str(t)
            # 흔한 모지바케 복구 시도
            try:
                t = t.encode("latin1").decode("utf-8")
            except Exception:
                pass
            if _is_readable(t):
                return t.strip()
    except Exception:
        pass
    # 첫 페이지 휴리스틱은 그대로 두고, 그래도 읽기 불가면 파일명 사용
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            if pdf.pages:
                text = pdf.pages[0].extract_text() or ""
                for ln in (l.strip() for l in text.splitlines()):
                    if len(ln) >= 3 and _is_readable(ln):
                        return ln[:120]
    except Exception:
        pass
    return path.stem

# [ADD OCR] OCR 전처리/실행
_HANGUL_ENG_NUM = re.compile(r"[ㄱ-ㅎ가-힣A-Za-z0-9]")

def _looks_gibberish(text: str) -> bool:
    """[OCR] 한/영/숫자 비율이 너무 낮거나 제어/기호가 과도하면 '깨짐'으로 간주"""
    if not text or len(text) < 30:
        return True
    good = len(_HANGUL_ENG_NUM.findall(text))
    ratio = good / max(1, len(text))
    return ratio < 0.25  # 경험치: 25% 미만이면 OCR 필요 가능성 높음

def _ocr_page(pil_img: Image.Image, lang: str = "kor+eng") -> str:
    # [ADD] OpenCV가 없으면 PIL+Tesseract만 빠르게 시도
    if not _CV_OK or not _TESSERACT_OK:
        return (pytesseract.image_to_string(pil_img, lang=lang) if _TESSERACT_OK else "").strip()
    """[OCR] 고해상도 + 이진화/잡음제거 후 Tesseract"""
    # 고해상도 확보
    img = np.array(pil_img.convert("RGB"))
    # 그레이스케일
    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    # 가우시안 블러로 노이즈 축소 후 OTSU 이진화
    gray = cv2.GaussianBlur(gray, (3,3), 0)
    _, binarized = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    # 약한 팽창/침식으로 글자 연결성 보정
    kernel = np.ones((2,2), np.uint8)
    proc = cv2.morphologyEx(binarized, cv2.MORPH_OPEN, kernel, iterations=1)
    # OCR
    config = "--oem 1 --psm 6 -c preserve_interword_spaces=1"
    proc = cv2.resize(proc, None, fx=1.5, fy=1.5, interpolation=cv2.INTER_CUBIC)
    text = pytesseract.image_to_string(proc, lang=lang, config=config)
    return text.strip()

# >>> FIX: OCR 전용 로더도 의존성 없으면 기본 파서로 강등
def load_pdf_ocr(path: Path, dpi: int = 400, lang: str = "kor+eng") -> List[Document]:
    """[OCR] 스캔/깨짐 PDF 전용 로더"""
    if not (_PDF2IMAGE_OK and _TESSERACT_OK):
        # ← OCR 불가: 기본 텍스트 파서 결과라도 리턴해 앱이 죽지 않게
        return _dedup(_build_docs_from_page_texts([], path) + load_pdf(path))

    pages = convert_from_path(str(path), dpi=dpi)   # poppler-utils 필요
    page_texts: List[str] = []
    docs: List[Document] = []

    for i, im in enumerate(pages, start=1):
        txt = _ocr_page(im, lang=lang)
        page_texts.append(txt or "")
        if txt:
            docs.append(Document(
                page_content=txt,
                metadata={"source": str(path), "type": "pdf", "page": i, "kind": "ocr_text"}
            ))

    # 제목/요약 메타 보강 (이하 기존 그대로)
    if docs:
        head = " ".join((docs[0].page_content or "").split()[:60])
        docs.insert(0, Document(page_content=f"[TITLE] {path.stem}", metadata={"source": str(path), "kind": "title"}))
        if head:
            docs.insert(1, Document(page_content=f"[SUMMARY] {path.name} 개요: {head}", metadata={"source": str(path), "kind": "summary"}))

    joined = _normalize_articles("\n".join(t for t in page_texts if t))
    section_docs = _build_section_docs_from_text(joined, path)
    docs = section_docs + docs

    return _dedup(docs)

    
def _preprocess_for_ocr(img: Image.Image) -> Image.Image:
    im = np.array(img.convert("RGB"))
    im = cv2.cvtColor(im, cv2.COLOR_RGB2GRAY)
    im = cv2.threshold(im, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
    return Image.fromarray(im)

# >>> OCR 스택이 없으면 빈 리스트를 반환하여 상위 로직이 텍스트 파서로 자연 강등
def _ocr_pdf_to_texts(pdf_path: Path, dpi: int = DEFAULT_OCR_DPI, lang: str = DEFAULT_OCR_LANG) -> List[str]:
    if not (_PDF2IMAGE_OK and _TESSERACT_OK):
        return []  # ← 의존성 없으면 OCR 시도 자체를 건너뜀

    images = convert_from_path(str(pdf_path), dpi=dpi)
    texts: List[str] = []
    for pil_img in images:
        proc = _preprocess_for_ocr(pil_img)
        txt = pytesseract.image_to_string(proc, lang=lang) if _TESSERACT_OK else ""
        texts.append((txt or "").strip())
    return texts


def _build_docs_from_page_texts(page_texts: List[str], path: Path) -> List[Document]:
    docs: List[Document] = []
    for idx, page_text in enumerate(page_texts, start=1):
        if page_text.strip():
            docs.append(Document(
                page_content=page_text,
                metadata={"source": str(path), "type": "pdf", "page": idx, "kind": "ocr_text"}
            ))
    # 간단 SUMMARY/TITLE 메타
    head_lines = []
    if page_texts:
        head_lines = [ln.strip() for ln in (page_texts[0].splitlines() if page_texts[0] else []) if ln.strip()][:10]
    head = " ".join(head_lines)[:400]
    if head:
        docs.append(Document(
            page_content=f"[SUMMARY] {path.name} 개요: {head}",
            metadata={"source": str(path), "kind": "summary"},
        ))
    title = _extract_pdf_title(path) or path.stem
    if title:
        docs.insert(0, Document(
            page_content=f"[TITLE] {title}",
            metadata={"source": str(path), "kind": "title"},
        ))
    # >>> '전체 본문'을 대상으로 범용 섹션 분해
    joined = "\n".join(t for t in page_texts if t)
    section_docs = _build_section_docs_from_text(joined, path)
    return _dedup(section_docs + docs)   # 섹션 문서를 우선 넣고 dedup

def _load_pdf_auto(path: Path) -> List[Document]:
    """표→문장 + 텍스트 스니펫을 함께 만들고,
       표 문장이 적으면 기본 파서와 병합. (텍스트가 너무 적으면 OCR 폴백)"""
    try:
        import pdfplumber
    except Exception:
        # pdfplumber 없으면 기본 텍스트만
        base = load_pdf(path)
        # [MOD OCR] 텍스트가 거의 없으면 OCR 폴백
        if sum(len(d.page_content) for d in base) < AUTO_OCR_MIN_CHARS:
            ocr_docs = _build_docs_from_page_texts(_ocr_pdf_to_texts(path), path)
            return _dedup(ocr_docs)
        return base

    docs: List[Document] = []
    first_text_lines: list[str] = []
    table_sentence_cnt = 0

    with pdfplumber.open(str(path)) as pdf:
        for pi, page in enumerate(pdf.pages, start=1):
            # 1) 표 → 문장
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for t in tables:
                sents = _table_to_sentences(t)
                table_sentence_cnt += len(sents)
                for s in sents:
                    docs.append(Document(
                        page_content=f"[SOURCE:{path.name}] {s}",
                        metadata={"source": str(path), "page": pi, "kind": "table_row"},
                    ))
            # 2) 텍스트 스니펫
            text = page.extract_text() or ""
            lines = [_norm(l) for l in text.splitlines() if _norm(l)]
            if pi == 1:
                first_text_lines = lines[:10]
            if lines:
                snippet = " ".join(lines[:20])[:1200]
                if snippet:
                    docs.append(Document(
                        page_content=f"[SOURCE:{path.name}] {snippet}",
                        metadata={"source": str(path), "page": pi, "kind": "page_text"},
                    ))

    # 3) 표가 거의 없으면 기본 파서와 병합
    threshold = settings.PDF_TABLE_THRESHOLD
    if table_sentence_cnt < threshold:
        legacy = load_pdf(path)
        docs = _dedup(legacy + docs)

    # 3.5) OCR 자동 폴백 판단 (요약/제목 추가 전!)
    #   - 후보 텍스트: 본문 성격만(기본 chunk/페이지 텍스트/표 문장)
    SRC_RE = re.compile(r'^\[SOURCE:.*?\]\s*')

    text_candidates = [
        SRC_RE.sub("", d.page_content)
        for d in docs
        if (d.metadata or {}).get("kind") in (None, "page_text", "table_row")
    ]
    full_text   = " ".join(text_candidates)
    total_chars = sum(len(t) for t in text_candidates)

    # 조건: 본문이 거의 없거나(길이), 있긴 한데 '깨짐'(가독성 낮음)
    if not text_candidates or total_chars < AUTO_OCR_MIN_CHARS or _looks_gibberish(full_text):
        # >>> FIX: OCR 스택이 있을 때만 OCR, 없으면 기본 파서로 강등
        if _PDF2IMAGE_OK and _TESSERACT_OK:
            ocr_docs = load_pdf_ocr(path, dpi=450, lang=DEFAULT_OCR_LANG)  # 고품질 파이프라인
            return _dedup(ocr_docs)
        return _dedup(load_pdf(path))  # ← 의존성 없으면 여기로
    
    joined_for_sections = _normalize_articles("\n".join(text_candidates))
    docs += _build_section_docs_from_text(joined_for_sections, path)
    docs = _dedup(docs)
    
    # 4) 요약/제목 메타 청크 (OCR 폴백 안 탔을 때만 진행)
    head = " ".join(first_text_lines)[:400]
    if head:
        docs.append(Document(
            page_content=f"[SUMMARY] {path.name} 개요: {head}",
            metadata={"source": str(path), "kind": "summary"},
        ))
    title = _extract_pdf_title(path)
    if title:
        docs.insert(0, Document(
            page_content=f"[TITLE] {title}",
            metadata={"source": str(path), "kind": "title"},
        ))
    return docs


# -----------------------------
# 디스패처
# -----------------------------
def load_docs_any(
    path: Union[str, Path],
    chunk_size: int = 800,
    chunk_overlap: int = 120,
    parser: str = "auto",
) -> List[Document]:
    p = Path(path)
    ext = p.suffix.lower()

    # [CHANGE OCR] 강제 OCR 파서: 고품질 파이프라인 사용
    if parser == "ocr_hi":
        if ext != ".pdf":
            raise ValueError("ocr_hi 파서는 PDF에만 사용할 수 있습니다.")
        docs = load_pdf_ocr(p, dpi=450, lang=DEFAULT_OCR_LANG)  # 고해상도 + 강한 전처리
        return _chunk(_dedup(docs), chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    
    if ext == ".pdf":
        docs = _load_pdf_auto(p)  # auto: 표-친화 시도 후 부족하면 기본 병합, 너무 적으면 OCR 폴백
        return _chunk(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    if ext == ".csv":
        docs = load_csv(p)
        return _chunk(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    if ext == ".hwp":
        raise ValueError("HWP는 지원하지 않습니다.")
    if ext == ".pptx":
        docs = load_pptx(p)
    elif ext == ".xlsx":
        docs = load_xlsx(p)
    elif ext in (".txt", ".log", ".md"):
        docs = load_txt(p)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    return _chunk(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)